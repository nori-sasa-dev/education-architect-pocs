import io


def extract_text_from_pdf(file_bytes: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text)
    return "\n".join(texts)


def extract_text_from_excel(file_bytes: bytes) -> str:
    import pandas as pd
    dfs = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None, header=None)
    texts = []
    for sheet_name, df in dfs.items():
        texts.append(f"【シート: {sheet_name}】")
        texts.append(df.fillna("").to_string(index=False, header=False))
    return "\n".join(texts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    texts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            texts.append(f"【スライド {i}】")
            texts.extend(slide_texts)
    return "\n".join(texts)


def extract_text(uploaded_file) -> str:
    name = uploaded_file.name.lower()
    data = uploaded_file.read()
    if name.endswith(".pdf"):
        return extract_text_from_pdf(data)
    elif name.endswith((".xlsx", ".xls")):
        return extract_text_from_excel(data)
    elif name.endswith((".pptx", ".ppt")):
        return extract_text_from_pptx(data)
    return ""
